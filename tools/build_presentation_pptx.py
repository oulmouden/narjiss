# -*- coding: utf-8 -*-
"""Fabrique la presentation PowerPoint sonorisee a partir du dossier Word
« argumentaire direction » (FR ou AR).

  python tools/build_presentation_pptx.py --langue fr
  python tools/build_presentation_pptx.py --langue ar --voix sage

Chaine complete :
 1. lecture du .docx (texte + images, sans dependance externe : c'est du zip + xml) ;
 2. montage des diapositives avec python-pptx, charte du dossier
    (ocre BD6B16, encre 16211C, gris 6D766F) ;
 3. narration de chaque diapositive par la synthese vocale OpenAI (mp3 mis en
    cache dans docs/voix-presentation/<langue>/ : une relance ne repaie pas) ;
 4. insertion des sons et minutage des diapositives par PowerPoint lui-meme
    (COM), le son demarrant tout seul a l'arrivee sur la diapositive.

Le texte affiche est volontairement plus court que le texte lu : la voix dit le
paragraphe entier, l'ecran n'en garde que l'amorce.
"""

import argparse
import hashlib
import html
import os
import re
import shutil
import sys
import tempfile
import zipfile
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from pptx.util import Emu, Inches, Pt

RACINE = Path(__file__).resolve().parent.parent
DOCS = RACINE / "docs"

# ---------------------------------------------------------------- charte -----

ENCRE = RGBColor(0x16, 0x21, 0x1C)
OCRE = RGBColor(0xBD, 0x6B, 0x16)
GRIS = RGBColor(0x6D, 0x76, 0x6F)
CREME = RGBColor(0xFA, 0xF7, 0xF2)
BLANC = RGBColor(0xFF, 0xFF, 0xFF)
TRAIT = RGBColor(0xE2, 0xDC, 0xD2)

LARGEUR = Inches(13.333)
HAUTEUR = Inches(7.5)
MARGE = Inches(0.85)

POLICES = {
    "fr": {"titre": "Georgia", "texte": "Calibri"},
    "ar": {"titre": "Arial", "texte": "Arial"},
}

# ------------------------------------------------------- lecture du .docx ----

SOURCES = {
    "fr": "NarjissargumentairedirectionFR.docx",
    "ar": "NarjissargumentairedirectionAR.docx",
}


def lire_docx(chemin, dossier_media):
    """Retourne la liste des elements du document, dans l'ordre.

    Chaque element est {'t': style, 'x': texte} ou {'t': 'img', 'src': fichier}.
    Les deux versions FR et AR ont exactement la meme structure, index par
    index : le plan des diapositives plus bas s'appuie la-dessus.
    """
    with zipfile.ZipFile(chemin) as z:
        doc = z.read("word/document.xml").decode("utf-8")
        rels = z.read("word/_rels/document.xml.rels").decode("utf-8")
        medias = {}
        for nom in z.namelist():
            if nom.startswith("word/media/"):
                cible = dossier_media / Path(nom).name
                cible.write_bytes(z.read(nom))
                medias[Path(nom).name] = cible

    liens = dict(re.findall(r'Id="([^"]+)"[^>]*Target="media/([^"]+)"', rels))
    elements = []
    for p in re.findall(r"<w:p\b.*?</w:p>|<w:p\b[^>]*/>", doc, re.S):
        style = re.search(r'w:pStyle w:val="([^"]+)"', p)
        style = style.group(1) if style else "Normal"
        texte = "".join(re.findall(r"<w:t[^>]*>(.*?)</w:t>", p, re.S))
        embed = re.findall(r'r:embed="([^"]+)"', p)
        if embed:
            elements.append({"t": "img", "src": medias[liens[embed[0]]]})
        if texte.strip():
            elements.append({"t": style, "x": html.unescape(texte).strip()})
    return elements


# ------------------------------------------------ plan des diapositives ------
# Les nombres sont les index dans la liste renvoyee par lire_docx().

PLAN = [
    {"type": "titre", "kicker": 0, "titre": 1, "sous": 2, "date": 3,
     "image": 5, "legende": 6, "dit": [1, 2, 4, 6]},
    {"type": "texte", "h1": 7, "corps": [8, 9]},
    {"type": "puces", "h1": 7, "h2": 10, "puces": [11, 12, 13, 14, 15, 16]},
    {"type": "texte_img", "h1": 17, "h2": 18, "corps": [19, 20],
     "image": 21, "legende": 22},
    {"type": "texte", "h1": 17, "h2": 23, "corps": [24, 25]},
    {"type": "sommaire", "h1": 26, "etapes": [27, 31, 36, 40, 45]},
    {"type": "texte_img", "h1": 26, "h2": 27, "corps": [28],
     "image": 29, "legende": 30},
    {"type": "texte_img", "h1": 26, "h2": 31, "corps": [32, 33],
     "image": 34, "legende": 35},
    {"type": "texte_img", "h1": 26, "h2": 36, "corps": [37],
     "image": 38, "legende": 39},
    {"type": "texte_img", "h1": 26, "h2": 40, "corps": [41, 42],
     "image": 43, "legende": 44},
    {"type": "texte", "h1": 26, "h2": 45, "corps": [46]},
    {"type": "texte", "h1": 47, "h2": 48, "corps": [49, 50]},
    {"type": "texte_img", "h1": 47, "h2": 51, "corps": [52, 53],
     "image": 54, "legende": 55},
    {"type": "texte_img", "h1": 56, "h2": 57, "corps": [58],
     "image": 59, "legende": 60},
    {"type": "texte_img", "h1": 56, "h2": 61, "corps": [62, 63],
     "image": 64, "legende": 65},
    {"type": "texte_img", "h1": 56, "h2": 66, "corps": [67],
     "image": 68, "legende": 69},
    {"type": "texte", "h1": 70, "corps": [71]},
    {"type": "puces", "h1": 70, "h2": 72, "puces": [73, 74, 75, 76, 77]},
    {"type": "texte_img", "h1": 70, "h2": 78, "corps": [79],
     "image": 85, "legende": 86},
    {"type": "puces", "h1": 70, "h2": 80, "puces": [81, 82, 83, 84]},
    {"type": "tableau", "h1": 87, "corps": [88],
     "entetes": [89, 90, 91], "premiere_cellule": 92, "lignes": 12},
    {"type": "puces", "h1": 140, "h2": 141,
     "puces": [142, 143, 144, 145, 146]},
    {"type": "puces", "h1": 140, "h2": 147,
     "puces": [148, 149, 150, 151, 152]},
    {"type": "puces", "h1": 153, "h2": 154,
     "puces": [155, 156, 157, 158, 159, 160, 161]},
    {"type": "puces", "h1": 153, "h2": 162, "puces": [163, 164]},
    {"type": "puces", "h1": 153, "h2": 165, "puces": [166, 167, 168]},
    {"type": "puces", "h1": 169, "puces": [170, 171, 172]},
    {"type": "final", "corps": [173, 174]},
]

# ------------------------------------------------------------- narration ----

# Ce qui se prononce mal tel quel. Ne touche pas au texte affiche.
ORAL = {
    "fr": [
        ("narjissimmobiliere.com", "narjiss immobilière point com"),
        ("narjiss.company", "narjiss point company"),
        ("360°", "360 degrés"),
        ("°", " degrés"),
        ("«", ""), ("»", ""),
    ],
    "ar": [
        ("narjissimmobiliere.com", "نرجس إيموبيلير دوت كوم"),
        ("narjiss.company", "نرجس دوت كومباني"),
        ("360°", "360 درجة"),
        ("°", " درجة"),
        ("«", ""), ("»", ""),
    ],
}

CONSIGNE = {
    "fr": ("Voix de presentation devant la direction d'une entreprise : ton "
           "posé, chaleureux et assuré, debit calme, articulation nette, "
           "vraies pauses aux points. Francais de France, jamais precipite."),
    "ar": ("صوت تقديم أمام إدارة شركة: نبرة هادئة وواثقة ودافئة، إيقاع متمهّل، "
           "نطق واضح بالعربية الفصحى، ووقفات حقيقية عند النقاط. لا تسرع."),
}


def pour_la_voix(texte, langue):
    for avant, apres in ORAL[langue]:
        texte = texte.replace(avant, apres)
    return re.sub(r"\s+", " ", texte).strip()


def phrases(texte):
    """Decoupe en phrases. Le point de « narjiss.company » n'est pas suivi
    d'une espace : il ne coupe donc rien."""
    return [p for p in re.split(r"(?<=[.!?؟])\s+", texte) if p]


def amorce(texte, maxi):
    """Les premieres phrases entieres qui tiennent dans maxi caracteres."""
    sortie = ""
    for p in phrases(texte):
        if sortie and len(sortie) + len(p) > maxi:
            break
        sortie = (sortie + " " + p).strip()
    return sortie or texte[:maxi]


# ------------------------------------------------------ briques de mise en page

def _pPr(paragraphe):
    return paragraphe._p.get_or_add_pPr()


def aligner(paragraphe, rtl):
    """Sens de lecture : l'arabe se compose a droite, et PowerPoint a besoin
    de l'attribut rtl pour placer correctement ponctuation et chiffres."""
    pPr = _pPr(paragraphe)
    if rtl:
        pPr.set("rtl", "1")
        pPr.set("algn", "r")
    else:
        pPr.set("rtl", "0")


def puce_ocre(paragraphe, rtl):
    pPr = _pPr(paragraphe)
    pPr.set("marL", "236220")
    pPr.set("indent", "-236220")
    for tag, attrs in (("a:buClr", None), ("a:buSzPct", {"val": "90000"}),
                       ("a:buFont", {"typeface": "Arial"}),
                       ("a:buChar", {"char": "▪"})):
        el = pPr.makeelement(qn(tag), attrs or {})
        if tag == "a:buClr":
            couleur = el.makeelement(qn("a:srgbClr"), {"val": "BD6B16"})
            el.append(couleur)
        pPr.append(el)


def bloc(diapo, gauche, haut, large, haute):
    forme = diapo.shapes.add_textbox(gauche, haut, large, haute)
    cadre = forme.text_frame
    cadre.word_wrap = True
    cadre.margin_left = cadre.margin_right = 0
    cadre.margin_top = cadre.margin_bottom = 0
    return cadre


def ecrire(cadre, textes, *, taille, police, couleur, gras=False,
           italique=False, interligne=1.15, espace=Pt(0), rtl=False,
           puces=False, majuscules=False, espacement=None, premier=True):
    for i, texte in enumerate(textes):
        p = cadre.paragraphs[0] if (premier and i == 0) else cadre.add_paragraph()
        p.line_spacing = interligne
        if i:
            p.space_before = espace
        aligner(p, rtl)
        if puces:
            puce_ocre(p, rtl)
        r = p.add_run()
        r.text = texte.upper() if majuscules else texte
        r.font.size = taille
        r.font.name = police
        r.font.bold = gras
        r.font.italic = italique
        r.font.color.rgb = couleur
        if espacement:
            r.font._rPr.set("spc", str(int(espacement)))
        # l'arabe ne se contente pas de la police latine : PowerPoint choisit
        # la police « complex script » pour ces caracteres.
        if rtl:
            r.font._rPr.append(r.font._rPr.makeelement(
                qn("a:cs"), {"typeface": police}))
    return cadre


def rectangle(diapo, gauche, haut, large, haute, couleur):
    from pptx.enum.shapes import MSO_SHAPE
    forme = diapo.shapes.add_shape(MSO_SHAPE.RECTANGLE, gauche, haut,
                                   large, haute)
    forme.fill.solid()
    forme.fill.fore_color.rgb = couleur
    forme.line.fill.background()
    forme.shadow.inherit = False
    return forme


def image_ajustee(diapo, chemin, gauche, haut, large, haute, couvrir=False):
    """Insere l'image dans la boite sans la deformer : contenue par defaut,
    ou rognee pour remplir toute la boite si couvrir=True."""
    from PIL import Image  # Pillow est deja installe (chaine video)
    try:
        with Image.open(chemin) as im:
            ratio = im.width / im.height
    except Exception:
        ratio = 16 / 9
    if couvrir:
        forme = diapo.shapes.add_picture(
            str(chemin), Emu(int(gauche)), Emu(int(haut)),
            Emu(int(large)), Emu(int(haute)))
        boite = large / haute
        if boite > ratio:                       # trop haute : on rogne en haut
            rogne = (1 - ratio / boite) / 2
            forme.crop_top = forme.crop_bottom = rogne
        else:                                   # trop large : on rogne les cotes
            rogne = (1 - boite / ratio) / 2
            forme.crop_left = forme.crop_right = rogne
        return forme
    if large / haute > ratio:
        h, l = haute, int(haute * ratio)
    else:
        l, h = large, int(large / ratio)
    forme = diapo.shapes.add_picture(
        str(chemin), Emu(int(gauche + (large - l) / 2)),
        Emu(int(haut + (haute - h) / 2)), Emu(int(l)), Emu(int(h)))
    forme.line.color.rgb = TRAIT
    forme.line.width = Pt(0.75)
    return forme


def pied(diapo, numero, rtl, police):
    if numero is None:
        return
    y = HAUTEUR - Inches(0.62)
    rectangle(diapo, MARGE, y, LARGEUR - 2 * MARGE, Pt(0.75), TRAIT)
    c = bloc(diapo, MARGE, y + Inches(0.12), LARGEUR - 2 * MARGE, Inches(0.3))
    p = c.paragraphs[0]
    p.alignment = PP_ALIGN.RIGHT if not rtl else PP_ALIGN.LEFT
    r = p.add_run()
    r.text = "narjiss.company   ·   %d" % numero
    r.font.size = Pt(9)
    r.font.name = police
    r.font.color.rgb = GRIS


# ------------------------------------------------------------ les gabarits ---

def diapo_titre(prs, el, fiche, langue, rtl):
    d = prs.slides.add_slide(prs.slide_layouts[6])
    pol = POLICES[langue]
    demi = int(LARGEUR * 0.52)
    rectangle(d, 0, 0, LARGEUR, HAUTEUR, ENCRE)
    if rtl:
        image_ajustee(d, el[fiche["image"]]["src"], 0, 0,
                      LARGEUR - demi, HAUTEUR, couvrir=True)
        gx = LARGEUR - demi + Inches(0.2)
    else:
        image_ajustee(d, el[fiche["image"]]["src"], demi, 0,
                      LARGEUR - demi, HAUTEUR, couvrir=True)
        gx = Inches(0.95)
    largeur_texte = demi - Inches(1.5)

    c = bloc(d, gx, Inches(1.55), largeur_texte, Inches(0.4))
    ecrire(c, [el[fiche["kicker"]]["x"]], taille=Pt(13), police=pol["texte"],
           couleur=OCRE, gras=True, espacement=260, rtl=False)

    c = bloc(d, gx, Inches(2.25), largeur_texte, Inches(2.2))
    ecrire(c, [el[fiche["titre"]]["x"]], taille=Pt(38), police=pol["titre"],
           couleur=BLANC, interligne=1.12, rtl=rtl)

    rectangle(d, gx if not rtl else gx + largeur_texte - Inches(1.4),
              Inches(4.55), Inches(1.4), Pt(2.5), OCRE)

    c = bloc(d, gx, Inches(5.0), largeur_texte, Inches(1.2))
    ecrire(c, [el[fiche["sous"]]["x"], el[fiche["date"]]["x"]],
           taille=Pt(15), police=pol["texte"],
           couleur=RGBColor(0xC9, 0xCE, 0xC9), espace=Pt(8), rtl=rtl)
    return d


def diapo_final(prs, el, fiche, langue, rtl):
    d = prs.slides.add_slide(prs.slide_layouts[6])
    pol = POLICES[langue]
    rectangle(d, 0, 0, LARGEUR, HAUTEUR, ENCRE)
    c = bloc(d, MARGE, Inches(2.5), LARGEUR - 2 * MARGE, Inches(1.0))
    p = c.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = "narjiss.company"
    r.font.size = Pt(44)
    r.font.name = pol["titre"]
    r.font.color.rgb = BLANC
    rectangle(d, int((LARGEUR - Inches(1.4)) / 2), Inches(3.8),
              Inches(1.4), Pt(2.5), OCRE)
    c = bloc(d, Inches(2.4), Inches(4.3), LARGEUR - Inches(4.8), Inches(1.8))
    for i, idx in enumerate(fiche["corps"]):
        p = c.paragraphs[0] if i == 0 else c.add_paragraph()
        p.alignment = PP_ALIGN.CENTER
        p.line_spacing = 1.3
        if i:
            p.space_before = Pt(10)
        aligner(p, rtl)
        r = p.add_run()
        r.text = el[idx]["x"]
        r.font.size = Pt(14)
        r.font.name = pol["texte"]
        r.font.color.rgb = RGBColor(0xC9, 0xCE, 0xC9)
        if rtl:
            r.font._rPr.append(r.font._rPr.makeelement(
                qn("a:cs"), {"typeface": pol["texte"]}))
    return d


def entete(d, el, fiche, langue, rtl, largeur=None):
    """Surtitre (chapitre) + titre de la diapositive. Renvoie le bas occupe."""
    pol = POLICES[langue]
    largeur = largeur or (LARGEUR - 2 * MARGE)
    x = MARGE if not rtl else LARGEUR - MARGE - largeur
    y = Inches(0.75)
    titre_idx = fiche.get("h2", fiche.get("h1"))
    surtitre = fiche.get("h1") if "h2" in fiche else None

    if surtitre is not None:
        c = bloc(d, x, y, largeur, Inches(0.32))
        ecrire(c, [el[surtitre]["x"]], taille=Pt(11), police=pol["texte"],
               couleur=OCRE, gras=True, rtl=rtl,
               espacement=None if rtl else 180,
               majuscules=(langue == "fr"))
        y += Inches(0.46)

    c = bloc(d, x, y, largeur, Inches(1.0))
    ecrire(c, [el[titre_idx]["x"]], taille=Pt(27), police=pol["titre"],
           couleur=ENCRE, interligne=1.1, rtl=rtl)
    # l'arabe tient plus de signes par pouce que le latin en Georgia
    densite = 5.2 if rtl else 4.6
    signes_par_ligne = max(12, int(Emu(int(largeur)).inches * densite))
    lignes = max(1, -(-len(el[titre_idx]["x"]) // signes_par_ligne))
    y += Inches(0.5) * lignes + Inches(0.14)
    rectangle(d, x if not rtl else x + largeur - Inches(1.1), y,
              Inches(1.1), Pt(2.5), OCRE)
    return y + Inches(0.42)


def diapo_texte(prs, el, fiche, langue, rtl):
    d = prs.slides.add_slide(prs.slide_layouts[6])
    pol = POLICES[langue]
    rectangle(d, 0, 0, LARGEUR, HAUTEUR, BLANC)
    largeur = LARGEUR - 2 * MARGE - Inches(1.6)
    y = entete(d, el, fiche, langue, rtl, largeur)
    x = MARGE if not rtl else LARGEUR - MARGE - largeur
    textes = [el[i]["x"] for i in fiche["corps"]]
    if sum(len(t) for t in textes) > 620:
        textes = [amorce(t, 330) for t in textes]
    c = bloc(d, x, y, largeur, HAUTEUR - y - Inches(1.0))
    ecrire(c, textes, taille=Pt(16), police=pol["texte"], couleur=ENCRE,
           interligne=1.42, espace=Pt(12), rtl=rtl)
    return d


def diapo_texte_img(prs, el, fiche, langue, rtl):
    d = prs.slides.add_slide(prs.slide_layouts[6])
    pol = POLICES[langue]
    rectangle(d, 0, 0, LARGEUR, HAUTEUR, BLANC)
    colonne = Inches(4.9)
    y = entete(d, el, fiche, langue, rtl, colonne)
    x = MARGE if not rtl else LARGEUR - MARGE - colonne

    textes = [amorce(el[i]["x"], 250) for i in fiche["corps"][:2]]
    if sum(len(t) for t in textes) > 400:
        textes = textes[:1]
    c = bloc(d, x, y, colonne, HAUTEUR - y - Inches(1.0))
    ecrire(c, textes, taille=Pt(15), police=pol["texte"], couleur=ENCRE,
           interligne=1.42, espace=Pt(12), rtl=rtl)

    img_x = MARGE + colonne + Inches(0.7) if not rtl else MARGE
    img_l = LARGEUR - MARGE - colonne - Inches(0.7) - MARGE
    img_y = Inches(0.95)
    img_h = Inches(4.75)
    image_ajustee(d, el[fiche["image"]]["src"], img_x, img_y, img_l, img_h)
    c = bloc(d, img_x, img_y + img_h + Inches(0.22), img_l, Inches(0.9))
    ecrire(c, [el[fiche["legende"]]["x"]], taille=Pt(10.5),
           police=pol["texte"], couleur=GRIS, italique=(langue == "fr"),
           interligne=1.25, rtl=rtl)
    return d


def diapo_puces(prs, el, fiche, langue, rtl):
    d = prs.slides.add_slide(prs.slide_layouts[6])
    pol = POLICES[langue]
    rectangle(d, 0, 0, LARGEUR, HAUTEUR, BLANC)
    largeur = LARGEUR - 2 * MARGE - Inches(1.2)
    y = entete(d, el, fiche, langue, rtl, largeur)
    x = MARGE if not rtl else LARGEUR - MARGE - largeur
    items = [el[i]["x"] for i in fiche["puces"]]
    place = HAUTEUR - y - Inches(1.0)
    taille = Pt(16) if len(items) <= 5 else Pt(14)
    if sum(len(t) for t in items) > 900:
        taille = Pt(13)
    c = bloc(d, x, y, largeur, place)
    ecrire(c, items, taille=taille, police=pol["texte"], couleur=ENCRE,
           interligne=1.35, espace=Pt(13), rtl=rtl, puces=True)
    return d


def diapo_sommaire(prs, el, fiche, langue, rtl):
    d = prs.slides.add_slide(prs.slide_layouts[6])
    pol = POLICES[langue]
    rectangle(d, 0, 0, LARGEUR, HAUTEUR, CREME)
    largeur = LARGEUR - 2 * MARGE
    x = MARGE
    c = bloc(d, x, Inches(1.75), largeur, Inches(1.1))
    ecrire(c, [el[fiche["h1"]]["x"]], taille=Pt(34), police=pol["titre"],
           couleur=ENCRE, rtl=rtl)
    rectangle(d, x if not rtl else LARGEUR - MARGE - Inches(1.4),
              Inches(2.95), Inches(1.4), Pt(2.5), OCRE)
    etapes = [el[i]["x"] for i in fiche["etapes"]]
    colonne = int((largeur - Inches(0.6) * (len(etapes) - 1)) / len(etapes))
    ordre = etapes if not rtl else list(reversed(etapes))
    for i, texte in enumerate(ordre):
        cx = x + i * (colonne + Inches(0.6))
        rectangle(d, cx, Inches(3.9), colonne, Pt(2), OCRE)
        c = bloc(d, cx, Inches(4.15), colonne, Inches(1.6))
        ecrire(c, [texte], taille=Pt(15), police=pol["texte"], couleur=ENCRE,
               interligne=1.3, rtl=rtl)
    return d


def diapo_tableau(prs, el, fiche, langue, rtl):
    d = prs.slides.add_slide(prs.slide_layouts[6])
    pol = POLICES[langue]
    rectangle(d, 0, 0, LARGEUR, HAUTEUR, BLANC)
    y = entete(d, el, fiche, langue, rtl)
    c = bloc(d, MARGE if not rtl else MARGE, y,
             LARGEUR - 2 * MARGE, Inches(0.4))
    ecrire(c, [el[fiche["corps"][0]]["x"]], taille=Pt(13),
           police=pol["texte"], couleur=GRIS, rtl=rtl)
    y += Inches(0.55)

    entetes = [""] + [el[i]["x"] for i in fiche["entetes"]]
    depart = fiche["premiere_cellule"]
    lignes = [[el[depart + r * 4 + col]["x"] for col in range(4)]
              for r in range(fiche["lignes"])]
    if rtl:
        entetes = list(reversed(entetes))
        lignes = [list(reversed(l)) for l in lignes]

    large = LARGEUR - 2 * MARGE
    tab = d.shapes.add_table(len(lignes) + 1, 4, MARGE, y, large,
                             HAUTEUR - y - Inches(0.9)).table
    tab.first_row = False
    tab.horz_banding = False
    premiere = Inches(3.2)
    autres = int((large - premiere) / 3)
    for i in range(4):
        est_libelle = (i == 3) if rtl else (i == 0)
        tab.columns[i].width = premiere if est_libelle else Emu(autres)

    for j, texte in enumerate(entetes):
        cell = tab.cell(0, j)
        p = cell.text_frame.paragraphs[0]
        aligner(p, rtl)
        r = p.add_run()
        r.text = texte
        r.font.size = Pt(11.5)
        r.font.bold = True
        r.font.name = pol["texte"]
        r.font.color.rgb = BLANC
        cell.fill.solid()
        cell.fill.fore_color.rgb = ENCRE
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        cell.margin_left = cell.margin_right = Inches(0.09)

    for i, ligne in enumerate(lignes, start=1):
        for j, texte in enumerate(ligne):
            est_libelle = (j == 3) if rtl else (j == 0)
            cell = tab.cell(i, j)
            p = cell.text_frame.paragraphs[0]
            aligner(p, rtl)
            r = p.add_run()
            r.text = texte
            r.font.size = Pt(10.5)
            r.font.name = pol["texte"]
            r.font.bold = est_libelle
            derniere = (j == 0) if rtl else (j == 3)
            r.font.color.rgb = OCRE if derniere else ENCRE
            cell.fill.solid()
            cell.fill.fore_color.rgb = CREME if i % 2 else BLANC
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            cell.margin_left = cell.margin_right = Inches(0.09)
            cell.margin_top = cell.margin_bottom = Inches(0.03)
    return d


GABARITS = {
    "titre": diapo_titre, "texte": diapo_texte, "texte_img": diapo_texte_img,
    "puces": diapo_puces, "sommaire": diapo_sommaire,
    "tableau": diapo_tableau, "final": diapo_final,
}


# ------------------------------------------------ texte lu par diapositive ---

def narration(el, fiche, langue):
    sep = ". " if langue == "fr" else ". "
    morceaux = []
    if fiche["type"] == "titre":
        morceaux = [el[i]["x"] for i in fiche["dit"]]
    elif fiche["type"] == "sommaire":
        liste = ", ".join(el[i]["x"] for i in fiche["etapes"])
        morceaux = [el[fiche["h1"]]["x"] + (" : " if langue == "fr" else ": ")
                    + liste]
    elif fiche["type"] == "tableau":
        morceaux = [el[fiche.get("h1")]["x"], el[fiche["corps"][0]]["x"]]
        depart = fiche["premiere_cellule"]
        joint = ", " if langue == "fr" else "، "
        deux = " : " if langue == "fr" else ": "
        for r in range(fiche["lignes"]):
            cells = [el[depart + r * 4 + c]["x"] for c in range(4)]
            morceaux.append(cells[0] + deux + joint.join(cells[1:]))
    else:
        titre = fiche.get("h2", fiche.get("h1"))
        if titre is not None:
            morceaux.append(el[titre]["x"])
        for i in fiche.get("corps", []) + fiche.get("puces", []):
            morceaux.append(el[i]["x"])
        if fiche.get("legende") is not None:
            morceaux.append(el[fiche["legende"]]["x"])
    texte = sep.join(m.rstrip(" .") for m in morceaux if m) + "."
    return pour_la_voix(texte, langue)


def synthetiser(textes, langue, voix, modele, dossier):
    """Un mp3 par diapositive, mis en cache sur le contenu exact."""
    import openai
    cle = re.search(r"^OPENAI_API_KEY=(.+)$",
                    (RACINE / "api" / ".env").read_text(encoding="utf-8"),
                    re.M).group(1).strip().strip('"').strip("'")
    client = openai.OpenAI(api_key=cle)
    dossier.mkdir(parents=True, exist_ok=True)
    sons = []
    for i, texte in enumerate(textes, start=1):
        empreinte = hashlib.sha1(
            ("%s|%s|%s" % (modele, voix, texte)).encode("utf-8")).hexdigest()[:10]
        chemin = dossier / ("diapo-%02d-%s.mp3" % (i, empreinte))
        if not chemin.exists():
            for vieux in dossier.glob("diapo-%02d-*.mp3" % i):
                vieux.unlink()
            reponse = client.audio.speech.create(
                model=modele, voice=voix, input=texte,
                instructions=CONSIGNE[langue], response_format="mp3")
            chemin.write_bytes(reponse.content)
            print("   voix %2d/%d  %5d signes  %s"
                  % (i, len(textes), len(texte), chemin.name))
        else:
            print("   voix %2d/%d  (cache)   %s" % (i, len(textes), chemin.name))
        sons.append(chemin)
    return sons


# ------------------------------- insertion des sons par PowerPoint (COM) -----

MSO_MEDIA = 16                        # msoMedia
MSO_DECLENCHEUR_AVEC_PRECEDENTE = 2   # msoAnimTriggerWithPrevious


def sonoriser(pptx, sons, rtl, enchainement=True, marge=Inches(0.28)):
    """Ajoute un son par diapositive, demarrage automatique, icone masquee,
    et cale la duree d'affichage sur la duree du commentaire."""
    import win32com.client
    ppt = win32com.client.Dispatch("PowerPoint.Application")
    ppt.Visible = True
    pres = ppt.Presentations.Open(str(pptx), ReadOnly=0, Untitled=0,
                                  WithWindow=0)
    try:
        for i, son in enumerate(sons, start=1):
            diapo = pres.Slides(i)
            # relancer le script ne doit pas empiler les sons
            for k in range(diapo.Shapes.Count, 0, -1):
                if diapo.Shapes(k).Type == MSO_MEDIA:
                    diapo.Shapes(k).Delete()
            cote = Inches(0.3)
            x = LARGEUR - marge - cote if not rtl else marge
            y = HAUTEUR - marge - cote
            forme = diapo.Shapes.AddMediaObject2(
                str(son), 0, -1, Emu(int(x)).pt, Emu(int(y)).pt,
                Emu(int(cote)).pt, Emu(int(cote)).pt)
            reglages = forme.AnimationSettings.PlaySettings
            reglages.PlayOnEntry = -1
            reglages.HideWhileNotPlaying = -1
            reglages.StopAfterSlides = 1
            # PlayOnEntry seul laisse l'effet en « au clic » dans la chronologie :
            # on bascule le declencheur sur « avec la precedente », c'est-a-dire
            # au moment ou la diapositive apparait.
            seq = diapo.TimeLine.MainSequence
            for k in range(1, seq.Count + 1):
                effet = seq.Item(k)
                if effet.Shape.Name == forme.Name:
                    effet.Timing.TriggerType = MSO_DECLENCHEUR_AVEC_PRECEDENTE
                    effet.Timing.TriggerDelayTime = 0.0
            # hors de la diapositive : invisible a l'ecran comme a l'export,
            # mais toujours dans la chronologie, donc toujours joue.
            forme.Left = -60.0
            forme.Top = -60.0
            duree = forme.MediaFormat.Length / 1000.0
            transition = diapo.SlideShowTransition
            transition.AdvanceOnClick = -1
            if enchainement:
                transition.AdvanceOnTime = -1
                transition.AdvanceTime = round(duree + 1.2, 1)
            print("   son  %2d/%d  %5.1f s" % (i, len(sons), duree))
        pres.Save()
    finally:
        pres.Close()
    return True


# ------------------------------------------------------------------ main ----

def main():
    ap = argparse.ArgumentParser(description=__doc__)
    ap.add_argument("--langue", choices=("fr", "ar"), default="fr")
    ap.add_argument("--voix", default="onyx",
                    help="voix OpenAI : onyx, ash, sage, nova, alloy...")
    ap.add_argument("--modele", default="gpt-4o-mini-tts")
    ap.add_argument("--sans-audio", action="store_true",
                    help="monter les diapositives sans passer par la voix")
    ap.add_argument("--sans-enchainement", action="store_true",
                    help="le son demarre seul mais la diapositive attend un clic")
    ap.add_argument("--sortie", default=None)
    args = ap.parse_args()

    langue = args.langue
    rtl = langue == "ar"
    source = DOCS / SOURCES[langue]
    if not source.exists():
        sys.exit("Introuvable : %s" % source)

    media = Path(tempfile.mkdtemp(prefix="pptx-media-"))
    try:
        el = lire_docx(source, media)
        print("Dossier lu : %d elements, %d diapositives prevues"
              % (len(el), len(PLAN)))

        prs = Presentation()
        prs.slide_width = LARGEUR
        prs.slide_height = HAUTEUR
        textes = []
        for numero, fiche in enumerate(PLAN, start=1):
            diapo = GABARITS[fiche["type"]](prs, el, fiche, langue, rtl)
            if fiche["type"] not in ("titre", "final"):
                pied(diapo, numero, rtl, POLICES[langue]["texte"])
            dit = narration(el, fiche, langue)
            textes.append(dit)
            diapo.notes_slide.notes_text_frame.text = dit

        sortie = Path(args.sortie) if args.sortie else (
            DOCS / ("Narjiss-argumentaire-direction-%s.pptx" % langue.upper()))
        prs.save(str(sortie))
        signes = sum(len(t) for t in textes)
        print("Diapositives montees : %s (%d signes a dire)" % (sortie, signes))

        if args.sans_audio:
            return
        dossier_voix = DOCS / "voix-presentation" / langue
        print("Synthese vocale (%s, voix %s)" % (args.modele, args.voix))
        sons = synthetiser(textes, langue, args.voix, args.modele, dossier_voix)
        print("Insertion des sons par PowerPoint")
        sonoriser(sortie, sons, rtl, not args.sans_enchainement)
        print("Termine : %s" % sortie)
    finally:
        shutil.rmtree(media, ignore_errors=True)


if __name__ == "__main__":
    main()
